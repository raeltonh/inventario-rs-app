import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib as mpl
import math
from scipy.stats import norm
import altair as alt
# --- Data containers ---
consumption_rates = {}
initial_stock    = {}
color_prices     = {}

# --- Check for statsmodels ARIMA ---
try:
    from statsmodels.tsa.arima.model import ARIMA
    HAS_ARIMA = True
except ModuleNotFoundError:
    HAS_ARIMA = False

# --- Global Styling (Matplotlib) ---
mpl.rcParams.update({
    'font.size': 10,
    'axes.titlesize': 12,
    'axes.labelsize': 10,
    'xtick.labelsize': 8,
    'ytick.labelsize': 8,
    'legend.fontsize': 8,
    'figure.figsize': (4, 3)
})

# --- Page Configuration ---
st.set_page_config(page_title="Inventory Simulation - RS Policy", layout="wide")
st.title("📦 Inventory Simulation - RS Policy")

# --- Default Parameters ---
lead_time       = 4
review_period   = 8
service_level   = 0.95
unit_price      = 40.0
holding_rate    = 0.17

demand1_color   = 15000.0
demand2_color   = 12000.0
demand1_white   = 10000.0
demand2_white   = 8000.0
selected_colors = ["Cyan", "Magenta"]

# --- Sidebar Settings ---
unit_option = st.sidebar.radio("Demand unit:", ("m²", "linear m"), index=0)

st.sidebar.header("⚙️ Settings")
enabled_features = st.sidebar.multiselect(
    "Enable features:",
    ["Scenario Comparison", "Forecast (ARIMA)", "Parameter Optimization", "Cost Summary"],
    default=["Scenario Comparison", "Forecast (ARIMA)", "Parameter Optimization", "Cost Summary"]
)

all_colors = ["Cyan","Magenta","Yellow","Black","Green","Red","DuoSoft","White"]
selected_colors = st.sidebar.multiselect("Select ink colors:", all_colors, default=selected_colors)

if st.sidebar.button("🔄 Reset all inputs"):
    st.session_state.clear()
    st.experimental_rerun()

# --- Timeline & Validation ---
dates = ["Aug/25","Sep/25","Oct/25","Nov/25","Dec/25","Jan/26"]
weeks_per_month = 4

if (lead_time + review_period) > len(dates) * weeks_per_month:
    st.error(f"Lead time + review period ({lead_time + review_period} weeks) exceeds horizon ({len(dates)*weeks_per_month} weeks).")

# --- Sidebar Inputs Continued ---
start_month    = st.sidebar.selectbox("Order start month:", dates, key='start_month')
lead_time      = st.sidebar.number_input("Lead time (weeks)", 1, 52, lead_time, key='lead_time')
review_period  = st.sidebar.number_input("Review period R (weeks)", 1, 52, review_period, key='review_period')
service_level  = st.sidebar.slider("Service level (%)", 80, 99, int(service_level*100), key='service_level') / 100
unit_price     = st.sidebar.number_input("Unit price ($/L)", 0.0, 1000.0, unit_price, key='unit_price')
holding_rate   = st.sidebar.slider("Annual holding cost rate (%)", 1, 50, int(holding_rate*100), key='holding_rate') / 100

label_unit = 'm²' if unit_option=='m²' else 'linear m'

demand1_color = st.sidebar.number_input(
    f"Demand ({label_unit}) Scenario 1 (colors)",
    0.0, 1e6, demand1_color,
    key='d1_color'
)
demand2_color = st.sidebar.number_input(
    f"Demand ({label_unit}) Scenario 2 (colors)",
    0.0, 1e6, demand2_color,
    key='d2_color'
)
demand1_white = st.sidebar.number_input(
    f"Demand ({label_unit}) Scenario 1 (white)",
    0.0, 1e6, demand1_white,
    key='d1_white'
)
demand2_white = st.sidebar.number_input(
    f"Demand ({label_unit}) Scenario 2 (white)",
    0.0, 1e6, demand2_white,
    key='d2_white'
)

# --- Simulation Function ---
def simulate(color, demand_m2):
    usage        = round((demand_m2 * consumption_rates[color]) / 1000, 2)
    std_month    = usage * 0.10
    total_lead   = lead_time + review_period
    demand_lead  = round((usage / weeks_per_month) * total_lead, 2)
    std_lead     = round(std_month * math.sqrt(total_lead / weeks_per_month), 2)
    z            = norm.ppf(service_level)
    safety_stock = round(z * std_lead, 2)

    Q = round(demand_lead + safety_stock, 2)
    S = math.ceil(Q / 5) * 5

    # Arrival & order timing
    step = math.floor(total_lead / weeks_per_month)
    idx = dates.index(start_month)
    arrival_months = []
    for _ in dates:
        idx += step
        if idx < len(dates):
            arrival_months.append(dates[idx])
    offset = math.ceil(lead_time / weeks_per_month)
    order_months = [
        dates[max(dates.index(m) - offset, 0)]
        for m in arrival_months
    ]

    # Build inventory table
    stock = initial_stock.get(color, 0)
    rows = []
    for m in dates:
        arrival     = S if m in arrival_months else 0
        start_stock = stock + arrival
        end_stock   = round(start_stock - usage, 2)
        status      = "Shortage" if end_stock < safety_stock else "OK"
        is_order    = m in order_months
        rows.append({
            "Month": m,
            "Start+Arrival": start_stock,
            "Usage": usage,
            "Arrival": arrival,
            "End Stock": end_stock,
            "Status": status,
            "Order": is_order
        })
        stock = end_stock

    df = pd.DataFrame(rows)
    df['Month'] = pd.Categorical(df['Month'], categories=dates, ordered=True)
    df = df.sort_values('Month')

    # Cost metrics
    cyc_cost      = (S/2) * unit_price * holding_rate
    saf_cost      = safety_stock * unit_price * holding_rate
    inv_value     = S * unit_price
    orders_per_yr = round(12 / (total_lead / weeks_per_month), 2)

    metrics = {
        'Usage': usage,
        'Safety Stock': safety_stock,
        'S': S,
        'Cycle Cost': cyc_cost,
        'Safety Cost': saf_cost,
        'Inventory Value': inv_value,
        'Orders/yr': orders_per_yr,
        'Order Months': order_months
    }
    return df, metrics

# --- Plot Function (Matplotlib) ---
def plot_df(df, color, safety_stock):
    cmap = {
        'Cyan':'#00AEEF','Magenta':'#FF00FF','Yellow':'#FFFF00',
        'Black':'#444444','Green':'#00CC66','Red':'#FF4444',
        'DuoSoft':'#FFA500','White':'#FFFFFF'
    }
    fig, ax = plt.subplots()
    ax.bar(df['Month'], df['End Stock'], color='lightgray', label='End Stock')
    ax.bar(
        df['Month'],
        df['Start+Arrival'] - df['End Stock'],
        bottom=df['End Stock'],
        color=cmap.get(color, '#888888'),
        label='Start+Arrival'
    )
    ax.axhline(
        safety_stock,
        linestyle='--',
        color=cmap.get(color, '#888888'),
        label='Safety Stock'
    )
    ax.set_xticks(range(len(dates)))
    ax.set_xticklabels(dates, rotation=45)
    ax.legend()
    return fig

# --- Inputs & Visualization ---
st.header("Color Data")
consumption_rates, initial_stock, color_prices = {}, {}, {}
cols = st.columns(3)
for i, color in enumerate(selected_colors):
    with cols[i % 3]:
        consumption_rates[color] = st.number_input(
            f"Consumption (ml/{label_unit}) - {color}", 0.0, 10.0, 3.45, key=f"cons_{color}"
        )
        _, m1 = simulate(color, demand1_white if color=='White' else demand1_color)
        _, m2 = simulate(color, demand2_white if color=='White' else demand2_color)
        st.markdown(
            f"**Order-up-to level (S) Scenario 1:** {m1['S']} L   "
            f"**Scenario 2:** {m2['S']} L"
        )
        initial_stock[color] = st.number_input(
            f"Initial stock (L) - {color}", 0.0, 10000.0, 0.0, key=f"stock_{color}"
        )
        color_prices[color] = st.number_input(
            f"Price ($/L) - {color}", 0.0, 1000.0, unit_price, key=f"price_{color}"
        )

# --- Cost Summary ---
if "Cost Summary" in enabled_features:
    st.header("Cost Summary")
    with st.expander("Summary & Formulas"):
        rows = []
        for color in selected_colors:
            _, m = simulate(color, demand1_white if color=='White' else demand1_color)
            rows.append({
                'Color': color,
                'Cycle Cost ($)': f"${m['Cycle Cost']:.2f}",
                'Safety Cost ($)': f"${m['Safety Cost']:.2f}",
                'Inv Value ($)': f"${m['Inventory Value']:.2f}",
                'Orders/yr': f"{m['Orders/yr']:.1f}",
                'Fill rate': f"{service_level*100:.0f}%"
            })
        cost_df = pd.DataFrame(rows)
        st.table(cost_df)
        csv = cost_df.to_csv(index=False).encode('utf-8')
        st.download_button('📥 Export Cost Summary CSV', data=csv,
                           file_name='cost_summary.csv', mime='text/csv')
        st.markdown("""
**Formulas:**
- Cycle Cost = (S/2) × unit_price × holding_rate  
- Safety Cost = safety_stock × unit_price × holding_rate  
- Inv Value  = S × unit_price  
- Orders/yr = 12 / ((lead_time + review_period)/4)  
- Fill rate = service_level  

**Legend:**  
- **Cycle Cost:** annual cost to hold average cycle inventory.  
- **Safety Cost:** annual cost to maintain safety stock.  
- **Inv Value:** total value of inventory at level S.  
- **Orders/yr:** estimated number of orders per year.  
- **Fill rate:** probability of meeting demand without stockouts.
""" )

# --- Scenario Comparison & Interactive Charts ---
cmap = {
    'Cyan':'#00AEEF','Magenta':'#FF00FF','Yellow':'#FFFF00',
    'Black':'#444444','Green':'#00CC66','Red':'#FF4444',
    'DuoSoft':'#FFA500','White':'#FFFFFF'
}

for color in selected_colors:
    df1, m1 = simulate(color, demand1_white if color=='White' else demand1_color)
    df2, m2 = simulate(color, demand2_white if color=='White' else demand2_color)

    # Always show shortages with recommended order timing
    if df1['Status'].eq('Shortage').any():
        shortage_months = df1.loc[df1['Status']=='Shortage', 'Month'].tolist()
        recommended_weeks = [f"first week of {m}" for m in m1['Order Months']]
        st.warning(
            f"⚠️ Shortage for {color} (Scenario 1) occurred in {', '.join(shortage_months)}. "
            f"Recommend placing order in the {', '.join(recommended_weeks)}."
        )
    if df2['Status'].eq('Shortage').any():
        shortage_months = df2.loc[df2['Status']=='Shortage', 'Month'].tolist()
        recommended_weeks = [f"first week of {m}" for m in m2['Order Months']]
        st.warning(
            f"⚠️ Shortage for {color} (Scenario 2) occurred in {', '.join(shortage_months)}. "
            f"Recommend placing order in the {', '.join(recommended_weeks)}."
        )

    if "Scenario Comparison" in enabled_features:
        st.subheader(f"{color} Scenarios")
        col1, col2 = st.columns(2)
        with col1:
            st.subheader("Scenario 1")
            st.dataframe(df1)
        with col2:
            st.subheader("Scenario 2")
            st.dataframe(df2)

        chart1 = (
            alt.Chart(df1)
               .mark_bar()
               .encode(
                   x=alt.X('Month:N', sort=dates, title='Month'),
                   y=alt.Y('End Stock:Q', title='End Stock (L)'),
                   color=alt.value(cmap[color]),
                   tooltip=[
                       alt.Tooltip('Month:N', title='Month'),
                       alt.Tooltip('Usage:Q', title='Usage (L)'),
                       alt.Tooltip('Arrival:Q', title='Arrival (L)'),
                       alt.Tooltip('End Stock:Q', title='End Stock (L)'),
                       alt.Tooltip('Status:N', title='Status')
                   ]
               )
               .interactive()
        )
        chart2 = (
            alt.Chart(df2)
               .mark_bar()
               .encode(
                   x=alt.X('Month:N', sort=dates, title='Month'),
                   y=alt.Y('End Stock:Q', title='End Stock (L)'),
                   color=alt.value(cmap[color]),
                   tooltip=[
                       alt.Tooltip('Month:N', title='Month'),
                       alt.Tooltip('Usage:Q', title='Usage (L)'),
                       alt.Tooltip('Arrival:Q', title='Arrival (L)'),
                       alt.Tooltip('End Stock:Q', title='End Stock (L)'),
                       alt.Tooltip('Status:N', title='Status')
                   ]
               )
               .interactive()
        )
        col1.altair_chart(chart1, use_container_width=True)
        col2.altair_chart(chart2, use_container_width=True)

# --- Forecast (ARIMA) ---
if "Forecast (ARIMA)" in enabled_features:
    if not HAS_ARIMA:
        st.warning("ARIMA models require statsmodels. Module not installed.")
    else:
        st.header("📈 ARIMA Forecast")
        date_index = pd.to_datetime(dates, format='%b/%y').to_period('M').to_timestamp()
        forecasts = []
        for color in selected_colors:
            hist_values = [((demand1_white if color=='White' else demand1_color)
                             * consumption_rates[color] / 1000) for _ in dates]
            ts = pd.Series(hist_values, index=date_index)
            model = ARIMA(ts, order=(1,1,1))
            res = model.fit()
            forecast = res.get_forecast(steps=3)
            last_date = date_index[-1]
            fc_index = [last_date + pd.DateOffset(months=i) for i in range(1,4)]
            fc_series = pd.Series(forecast.predicted_mean.values, index=fc_index)
            ci = forecast.conf_int()
            ci.index = fc_index
            forecasts.append((color, ts, fc_series, ci))

        for i in range(0, len(forecasts), 2):
            cols = st.columns(2)
            for j, (color, ts, fc_series, ci) in enumerate(forecasts[i:i+2]):
                with cols[j]:
                    st.subheader(f"{color} Usage Forecast")
                    fig, ax = plt.subplots()
                    ax.plot(ts.index, ts.values, marker='o', label='Historical')
                    ax.plot(fc_series.index, fc_series.values,
                            marker='o', linestyle='--', label='Forecast')
                    ax.fill_between(fc_series.index,
                                    ci.iloc[:,0], ci.iloc[:,1],
                                    alpha=0.2)
                    ax.set_title(f"{color} Usage (L/month)")
                    ax.set_xlabel('Date')
                    ax.legend()
                    fig.autofmt_xdate()
                    st.pyplot(fig)

# --- Parameter Optimization & Cost Alternatives ---
with st.expander("Parameter Optimization & Cost Alternatives"):
    if "Parameter Optimization" in enabled_features and selected_colors:
        R_vals = [4, 8, 12]
        SL_vals = [0.90, 0.95, 0.99]
        opt = []
        for R in R_vals:
            for SL in SL_vals:
                prevR, prevSL = review_period, service_level
                review_period, service_level = R, SL
                total_cost = 0
                for color in selected_colors:
                    _, m = simulate(color, demand1_white if color=='White' else demand1_color)
                    total_cost += m['Cycle Cost'] + m['Safety Cost']
                opt.append({'R (wk)': R, 'SL': f"{int(SL*100)}%", 'Total Cost ($)': total_cost})
                review_period, service_level = prevR, prevSL
        df_opt = pd.DataFrame(opt).sort_values('Total Cost ($)')
        st.table(df_opt)
        best = df_opt.iloc[0]
        st.success(
            f"Optimal review period: {best['R (wk)']} weeks "
            f"at SL {best['SL']}, minimizing annual cost to ${best['Total Cost ($)']:.2f}."
        )

# --- PowerPoint Report Generation ---
try:
    from pptx import Presentation
    from pptx.util import Inches
    HAS_PPTX = True
except ModuleNotFoundError:
    HAS_PPTX = False

if "Cost Summary" in enabled_features and HAS_PPTX:
    prs = Presentation()
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Inventory Simulation Report"
    slide.placeholders[1].text = "Generated from RS Policy App"

    # Cost Summary slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Cost Summary"
    rows = [['Color','Cycle Cost','Safety Cost','Inv Value','Orders/yr','Fill rate']]
    for color in selected_colors:
        _, m = simulate(color, demand1_white if color=='White' else demand1_color)
        rows.append([
            color,
            f"${m['Cycle Cost']:.2f}",
            f"${m['Safety Cost']:.2f}",
            f"${m['Inventory Value']:.2f}",
            f"{m['Orders/yr']:.1f}",
            f"{service_level*100:.0f}%"
        ])
    x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4)
    table = slide.shapes.add_table(len(rows), len(rows[0]), x, y, cx, cy).table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            table.cell(r,c).text = str(val)

    # Chart slides per color
    for color in selected_colors:
        df1, m1 = simulate(color, demand1_white if color=='White' else demand1_color)
        fig = plot_df(df1, color, m1['Safety Stock'])
        img1 = f"{color}_scenario1.png"
        fig.savefig(img1)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"{color} Scenario 1 Inventory"
        slide.shapes.add_picture(img1, Inches(1), Inches(1.5), width=Inches(8))
        plt.close(fig)

        df2, m2 = simulate(color, demand2_white if color=='White' else demand2_color)
        fig = plot_df(df2, color, m2['Safety Stock'])
        img2 = f"{color}_scenario2.png"
        fig.savefig(img2)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"{color} Scenario 2 Inventory"
        slide.shapes.add_picture(img2, Inches(1), Inches(1.5), width=Inches(8))
        plt.close(fig)

    pptx_path = 'inventory_report.pptx'
    prs.save(pptx_path)
    with open(pptx_path, 'rb') as f:
        st.download_button(
            label='📥 Download Full Report (PPTX)',
            data=f,
            file_name=pptx_path,
            mime='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
